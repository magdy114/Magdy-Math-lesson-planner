from __future__ import annotations

import base64
import io
import json
import logging
import os
import re
import shutil
import time
import uuid
import zipfile
from dataclasses import dataclass, asdict
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Tuple

from dotenv import load_dotenv
from flask import Flask, flash, jsonify, redirect, render_template, request, send_file, url_for
from werkzeug.utils import secure_filename
from docx import Document
from docx.shared import Pt
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

try:
    from openai import OpenAI
except Exception:  # pragma: no cover
    OpenAI = None

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover
    PdfReader = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None

from curriculum_document_builder import build_long, build_medium
from curriculum_pdf_builder import build_long_pdf, build_medium_pdf
from curriculum_extractors import candidate_topics, extract_text as extract_curriculum_text
from curriculum_models import LongPlan, MediumPlan
from curriculum_ai import generate_long, generate_medium, refine_topics

BASE_DIR = Path(__file__).resolve().parent
TEMPLATE_PATH = BASE_DIR / "assets" / "Lesson_Plan_Template_AY2026_2027.docx"
LIBRARY_DIR = BASE_DIR / "generated_lessons"
UPLOAD_DIR = BASE_DIR / "uploads"
LOG_DIR = BASE_DIR / "logs"
META_PATH = LIBRARY_DIR / "library.json"
USAGE_PATH = BASE_DIR / "usage.json"
CURRICULUM_ROOT = BASE_DIR / "generated_plans"
CURRICULUM_UPLOAD_DIR = CURRICULUM_ROOT / "uploads"
CURRICULUM_JOB_DIR = CURRICULUM_ROOT / "jobs"
CURRICULUM_EXPORT_DIR = CURRICULUM_ROOT / "exports"

load_dotenv(BASE_DIR / ".env")

LIBRARY_DIR.mkdir(exist_ok=True)
UPLOAD_DIR.mkdir(exist_ok=True)
LOG_DIR.mkdir(exist_ok=True)
for _folder in (CURRICULUM_UPLOAD_DIR, CURRICULUM_JOB_DIR, CURRICULUM_EXPORT_DIR):
    _folder.mkdir(parents=True, exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(message)s",
    handlers=[
        logging.FileHandler(LOG_DIR / "error_log.txt", encoding="utf-8"),
        logging.StreamHandler(),
    ],
)
logger = logging.getLogger("magdy_lesson_planner")

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 35 * 1024 * 1024
app.config["UPLOAD_FOLDER"] = str(UPLOAD_DIR)
app.secret_key = os.getenv("FLASK_SECRET", "magdy-lesson-planner-local-secret")

ALLOWED_EXTENSIONS = {"pdf", "docx", "txt", "pptx", "png", "jpg", "jpeg", "webp"}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "webp"}
REQUIRED_FIELDS = [
    "keywords", "sdg", "strategies", "intervention", "learning_outcomes", "differentiation",
    "success_criteria", "starter", "main", "teacher_led", "student_led", "plenary", "kpi",
    "resources", "identity", "competency", "curriculum"
]
MAX_FIELD_CHARS = {
    "keywords": 300,
    "sdg": 260,
    "strategies": 520,
    "intervention": 520,
    "learning_outcomes": 700,
    "differentiation": 620,
    "success_criteria": 620,
    "starter": 430,
    "main": 620,
    "teacher_led": 520,
    "student_led": 520,
    "plenary": 400,
    "kpi": 450,
    "resources": 360,
    "identity": 330,
    "competency": 350,
    "curriculum": 420,
}


@dataclass
class LessonInput:
    index: int
    teacher: str
    subject: str
    class_name: str
    periods: str
    language: str
    topic: str
    date: str
    notes: str = ""
    source_text: str = ""
    source_file_name: str = ""
    ai_mode: str = "auto"
    access_code: str = ""


# ----------------------------- Small utilities -----------------------------

def env_int(name: str, default: int) -> int:
    try:
        return int(os.getenv(name, str(default)))
    except Exception:
        return default


def fmt_date(date_value: str) -> str:
    if not date_value:
        return ""
    for fmt in ("%Y-%m-%d", "%d/%m/%Y", "%m/%d/%Y"):
        try:
            return datetime.strptime(date_value, fmt).strftime("%d/%m/%Y")
        except ValueError:
            continue
    return date_value


def safe_filename(text: str, fallback: str = "lesson_plan") -> str:
    text = (text or fallback).strip()
    text = re.sub(r"[\\/:*?\"<>|]+", "-", text)
    text = re.sub(r"\s+", "_", text)
    text = re.sub(r"_+", "_", text)
    return text[:90].strip("_.-") or fallback


def clean_text(text: str, max_chars: int | None = None) -> str:
    text = str(text or "")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if max_chars and len(text) > max_chars:
        text = text[: max_chars - 3].rstrip() + "..."
    return text


def file_ext(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def status_payload() -> Dict[str, Any]:
    return {
        "ai_enabled": bool(os.getenv("OPENAI_API_KEY")) and OpenAI is not None,
        "model": os.getenv("OPENAI_MODEL", "gpt-5.5"),
        "template_exists": TEMPLATE_PATH.exists(),
        "access_code_required": False,
    }


# ----------------------------- Library and limits ----------------------------

def load_json(path: Path, default: Any) -> Any:
    if not path.exists():
        return default
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return default


def save_json(path: Path, data: Any) -> None:
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")


def load_library() -> List[Dict[str, str]]:
    return load_json(META_PATH, [])


def save_library(items: List[Dict[str, str]]) -> None:
    save_json(META_PATH, items[:500])


def add_to_library(record: Dict[str, str]) -> None:
    items = load_library()
    items.insert(0, record)
    save_library(items)


def validate_access_code(code: str) -> Tuple[bool, str]:
    # Access code is disabled for the school-sharing version.
    # The API key stays protected on the server via .env; teachers do not need any login code.
    return True, ""


def check_usage_limit(user_key: str) -> Tuple[bool, str]:
    today = datetime.now().strftime("%Y-%m-%d")
    data = load_json(USAGE_PATH, {})
    if data.get("date") != today:
        data = {"date": today, "total": 0, "users": {}}
    total_limit = env_int("DAILY_TOTAL_LIMIT", 300)
    user_limit = env_int("DAILY_USER_LIMIT", 25)
    user_key = safe_filename(user_key or "anonymous", "anonymous")[:80]
    total = int(data.get("total", 0))
    user_count = int(data.get("users", {}).get(user_key, 0))
    if total >= total_limit:
        return False, "تم الوصول للحد اليومي العام. حاول غدًا أو زد الحد من ملف .env."
    if user_count >= user_limit:
        return False, "تم الوصول للحد اليومي لهذا المستخدم."
    data["total"] = total + 1
    data.setdefault("users", {})[user_key] = user_count + 1
    save_json(USAGE_PATH, data)
    return True, ""


# ----------------------------- DOCX styling helpers ----------------------------

def set_paragraph_bidi(paragraph, rtl: bool) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    existing = p_pr.find(qn("w:bidi"))
    if rtl:
        if existing is None:
            existing = OxmlElement("w:bidi")
            p_pr.append(existing)
        existing.set(qn("w:val"), "1")
    elif existing is not None:
        p_pr.remove(existing)


def set_run_font(run, font_name: str, size: float, bold: bool = False) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    r_pr = run._element.get_or_add_rPr()
    r_fonts = r_pr.rFonts
    if r_fonts is None:
        r_fonts = OxmlElement("w:rFonts")
        r_pr.append(r_fonts)
    for key in ("w:ascii", "w:hAnsi", "w:cs"):
        r_fonts.set(qn(key), font_name)


def set_cell_text(cell, text: str, lang: str = "en", size: float = 8.0, bold: bool = False) -> None:
    rtl = lang == "ar"
    font = "Arial" if rtl else "Times New Roman"
    cell.text = ""
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.TOP
    text = clean_text(text)
    lines = text.split("\n") if text else [""]
    for i, line in enumerate(lines):
        paragraph = cell.paragraphs[0] if i == 0 else cell.add_paragraph()
        paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT if rtl else WD_ALIGN_PARAGRAPH.LEFT
        set_paragraph_bidi(paragraph, rtl)
        paragraph.paragraph_format.space_after = Pt(0)
        paragraph.paragraph_format.space_before = Pt(0)
        run = paragraph.add_run(line)
        set_run_font(run, font, size, bold=bold)


def fill_check(cell, lang: str = "en") -> None:
    set_cell_text(cell, "✓", lang=lang, size=12, bold=True)


# ----------------------------- Upload extraction -----------------------------

def extract_docx_text(path: Path, limit: int = 6000) -> str:
    try:
        doc = Document(str(path))
        pieces: List[str] = []
        for p in doc.paragraphs:
            if p.text.strip():
                pieces.append(p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                line = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if line:
                    pieces.append(line)
        return clean_text("\n".join(pieces), limit)
    except Exception as exc:
        logger.exception("DOCX extraction failed")
        return f"[Could not extract DOCX text: {exc}]"


def extract_pdf_text(path: Path, limit: int = 6000) -> str:
    if PdfReader is None:
        return "[PDF extraction package is not installed.]"
    try:
        reader = PdfReader(str(path))
        pieces: List[str] = []
        for page in reader.pages[:10]:
            txt = page.extract_text() or ""
            if txt.strip():
                pieces.append(txt.strip())
        return clean_text("\n".join(pieces), limit)
    except Exception as exc:
        logger.exception("PDF extraction failed")
        return f"[Could not extract PDF text: {exc}]"


def extract_pptx_text(path: Path, limit: int = 6000) -> str:
    if Presentation is None:
        return "[PowerPoint extraction package is not installed.]"
    try:
        prs = Presentation(str(path))
        pieces: List[str] = []
        for i, slide in enumerate(prs.slides[:12], start=1):
            slide_lines: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
            if slide_lines:
                pieces.append(f"Slide {i}: " + " | ".join(slide_lines))
        return clean_text("\n".join(pieces), limit)
    except Exception as exc:
        logger.exception("PPTX extraction failed")
        return f"[Could not extract PPTX text: {exc}]"


def extract_txt_text(path: Path, limit: int = 6000) -> str:
    for enc in ("utf-8-sig", "utf-8", "cp1256", "cp1252"):
        try:
            return clean_text(path.read_text(encoding=enc, errors="ignore"), limit)
        except Exception:
            continue
    return ""


def save_uploaded_file(storage, prefix: str) -> Tuple[str, str]:
    if not storage or not storage.filename:
        return "", ""
    ext = file_ext(storage.filename)
    if ext not in ALLOWED_EXTENSIONS:
        return storage.filename, f"[Unsupported file type: {ext}]"
    filename = secure_filename(storage.filename)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    saved = UPLOAD_DIR / f"{prefix}_{stamp}_{filename}"
    storage.save(str(saved))
    if ext == "docx":
        text = extract_docx_text(saved)
    elif ext == "pdf":
        text = extract_pdf_text(saved)
    elif ext == "pptx":
        text = extract_pptx_text(saved)
    elif ext == "txt":
        text = extract_txt_text(saved)
    else:
        text = ""
    return storage.filename, text


def save_pasted_image(data_url: str, prefix: str) -> str:
    if not data_url or "," not in data_url:
        return ""
    try:
        header, payload = data_url.split(",", 1)
        ext = "jpg" if "jpeg" in header or "jpg" in header else "png"
        raw = base64.b64decode(payload)
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        path = UPLOAD_DIR / f"{prefix}_{stamp}_pasted_image.{ext}"
        path.write_bytes(raw)
        return path.name
    except Exception:
        logger.exception("Could not save pasted image")
        return ""


# ----------------------------- Content generation -----------------------------

def topic_family(topic: str, source_text: str = "") -> str:
    t = f"{topic or ''} {source_text or ''}".lower()
    if any(w in t for w in ["مشتق", "اشتق", "derivative", "differentiation", "tangent", "مماس", "rate of change"]):
        return "derivatives"
    if any(w in t for w in ["نها", "limit", "continuity", "اتصال", "approach"]):
        return "limits"
    if any(w in t for w in ["تكامل", "integral", "area under", "definite", "antiderivative", "arc length", "طول المنحنى"]):
        return "integrals"
    if any(w in t for w in ["لوغ", "log", "exponential", "أسي", "أسية", "ln", "growth", "decay"]):
        return "logs"
    if any(w in t for w in ["مثلث", "trig", "sine", "cos", "tan", "جيب", "جا", "جتا", "radian", "unit circle"]):
        return "trig"
    if any(w in t for w in ["دالة", "دوال", "function", "inverse", "composition", "domain", "range", "asymptote", "تقارب"]):
        return "functions"
    if any(w in t for w in ["matrix", "مصفوف", "determinant", "cramer"]):
        return "matrices"
    if any(w in t for w in ["vector", "متجه", "dot product", "cross product"]):
        return "vectors"
    return "general"


def source_keywords(source_text: str, lang: str) -> str:
    if not source_text:
        return ""
    terms = []
    candidates = [
        "derivative", "tangent", "rate of change", "limit", "continuity", "integral", "area", "logarithm", "exponential",
        "function", "domain", "range", "radian", "unit circle", "matrix", "vector", "slope", "asymptote", "arc length",
        "المشتقة", "المماس", "معدل التغير", "النهاية", "الاتصال", "التكامل", "المساحة", "اللوغاريتم", "الدالة", "المجال", "المدى", "الراديان", "المصفوفات", "المتجهات", "خط التقارب", "طول المنحنى"
    ]
    low = source_text.lower()
    for c in candidates:
        if c.lower() in low and c not in terms:
            terms.append(c)
    return ("، " if lang == "ar" else ", ").join(terms[:8])


def offline_content(lesson: LessonInput) -> Dict[str, str]:
    lang = lesson.language
    fam = topic_family(lesson.topic, lesson.source_text)
    topic = lesson.topic.strip() or ("الدرس" if lang == "ar" else "the lesson")
    subject = lesson.subject.strip() or ("رياضيات" if lang == "ar" else "Mathematics")
    class_name = lesson.class_name.strip() or ("الثاني عشر متقدم" if lang == "ar" else "Grade 12 Advanced")
    sk = source_keywords(lesson.source_text, lang)
    note = lesson.notes.strip()

    ar_bank = {
        "derivatives": ["المشتقة، ميل المماس، معدل التغير اللحظي، قاعدة القوى، التمثيل البياني", "تفسير المشتقة كميل للمماس ومعدل تغير لحظي ثم استخدامها في مسائل تطبيقية.", "سؤال استرجاع: احسب ميل قاطع ثم ناقش ماذا يحدث عند اقتراب النقطتين.", "نمذجة حل لإيجاد المشتقة ثم تدريب موجه ثم تحدي في معدل التغير.", "الخلط بين قيمة الدالة وقيمة المشتقة."],
        "limits": ["النهاية، الاقتراب، الاتصال، عدم التعيين، التحليل، السلوك البياني", "تقدير النهايات بيانياً وعددياً وجبرياً وربطها بالاتصال.", "جدول قيم حول نقطة وسؤال: ما القيمة التي تقترب منها الدالة؟", "نموذج تعويض مباشر ثم تحليل عند الحاجة، تدريب من رسم وجدول، وتحدي في حالة ثقب أو خط تقارب.", "الاعتقاد أن النهاية تساوي دائماً قيمة الدالة."],
        "integrals": ["التكامل، المساحة، التراكم، الدالة الأصلية، النظرية الأساسية", "فهم التكامل كمساحة/تراكم واستخدامه لحل مسائل رياضية وسياقية.", "تقدير مساحة تحت منحنى باستخدام مستطيلات قصيرة.", "شرح بصري للتراكم، تدريب على القاعدة المناسبة، وتطبيق سياقي في السرعة أو الاستهلاك.", "الخلط بين المساحة الهندسية والتكامل المحدد ذي الإشارة."],
        "logs": ["الدوال الأسية، اللوغاريتمات، خصائص اللوغاريتم، النمو والانحلال", "استخدام خصائص اللوغاريتمات والدوال الأسية في التبسيط وحل المعادلات.", "مطابقة بين الصورة الأسية والصورة اللوغاريتمية.", "تحويل بين الصورتين، تدريب على الخصائص، ومسألة نمو/انحلال.", "تطبيق log(a+b)=log a + log b خطأ."],
        "trig": ["الدوال المثلثية، الدائرة المثلثية، الزوايا، الراديان، التحويلات", "ربط الزوايا والدائرة المثلثية والتمثيل البياني لحل مسائل دقيقة.", "مراجعة قيم الزوايا الخاصة وتحديد الإشارات حسب الربع.", "شرح بصري، تدريب على القيم والتحويل بين الدرجات والراديان، ونشاط زوجي في السعة والدورة.", "الخلط بين الدرجات والراديان."],
        "functions": ["الدوال، المجال والمدى، التركيب، العكس، التحويلات، التمثيل البياني", "تحليل خصائص الدوال وتمثيلها وتفسير التحويلات أو التركيب/العكس.", "سؤال سريع عن المجال والمدى من رسم بياني.", "نموذج تحليل خصائص دالة، تدريب موجه، ومهمة تربط الجبر بالرسم.", "الخلط بين المجال والمدى أو ترتيب التحويلات."],
        "matrices": ["المصفوفات، المحدد، المعكوس، العمليات، الأنظمة الخطية", "استخدام المصفوفات والعمليات المناسبة لحل مسائل جبرية وسياقية.", "قراءة عنصر من مصفوفة وتفسير معنى الصف والعمود.", "نموذج محلول للعملية، تدريب موجه، وتحدي للتحقق بالآلة الحاسبة.", "إهمال شرط الأبعاد في ضرب المصفوفات."],
        "vectors": ["المتجهات، المقدار، الاتجاه، المركبات، الضرب النقطي", "تمثيل المتجهات وتحليلها واستخدامها في تفسير الحركة والعلاقات الهندسية.", "عرض سهمين على شبكة وتحديد المركبات والمقدار.", "تحويل المتجه لمركبات، تدريب على العمليات، وتطبيق في الحركة أو القوى.", "الخلط بين المقدار والاتجاه."],
        "general": [f"{topic}، المفاهيم الأساسية، التبرير، حل المشكلات، التواصل الرياضي", f"بناء فهم عميق لدرس {topic} من خلال الشرح الموجه والتدريب المتدرج.", f"سؤال استرجاع مرتبط بمتطلبات درس {topic}.", f"نموذج محلول، تدريب موجه، ومهمة تفكير عليا في {topic}.", "خطأ شائع يتم تحديده من إجابات البداية."],
    }
    en_bank = {
        "derivatives": ["Derivative, tangent slope, instantaneous rate of change, power rule, graph interpretation", "Interpret the derivative as tangent slope and instantaneous rate of change, then apply it to real problems.", "Retrieval: find a secant slope, then discuss what happens as the points get closer.", "Model a derivative solution, guided practice, then a rate-of-change challenge.", "Confusing function value with derivative value."],
        "limits": ["Limit, approaching, continuity, indeterminate form, factorisation, graph behaviour", "Estimate limits graphically, numerically, and algebraically, and link them to continuity.", "Show a table near a value and ask: what value is the function approaching?", "Direct substitution, factorisation when needed, graph/table practice, and a challenge with holes or asymptotes.", "Assuming the limit always equals the function value."],
        "integrals": ["Integral, area, accumulation, antiderivative, Fundamental Theorem", "Understand integration as area/accumulation and use it in mathematical and contextual problems.", "Estimate area under a curve using rectangles.", "Visual accumulation model, rule selection practice, and contextual application in velocity or consumption.", "Confusing geometric area with signed definite integral."],
        "logs": ["Exponential functions, logarithms, log properties, growth and decay", "Use logarithmic and exponential properties to simplify and solve equations.", "Match exponential and logarithmic forms.", "Convert between forms, practise log properties, and solve a growth/decay problem.", "Using log(a+b)=log a + log b incorrectly."],
        "trig": ["Trigonometric functions, unit circle, angles, radians, transformations", "Connect angles, the unit circle, and graphs of trigonometric functions to solve accurate problems.", "Review special angle values and signs by quadrant.", "Visual model, values and degree/radian conversion, paired activity on amplitude and period.", "Confusing degrees and radians."],
        "functions": ["Functions, domain and range, composition, inverse, transformations, graphs", "Analyse function properties and interpret transformations, composition, or inverse relationships.", "Quick question on domain and range from a graph.", "Model function analysis, guided practice, and a task linking algebra and graphing.", "Confusing domain/range or transformation order."],
        "matrices": ["Matrices, determinant, inverse, operations, linear systems", "Use matrices and suitable operations to solve algebraic and contextual problems.", "Read a matrix element and interpret row/column meaning.", "Worked example, guided practice, and calculator verification challenge.", "Ignoring dimension rules in matrix multiplication."],
        "vectors": ["Vectors, magnitude, direction, components, dot product", "Represent and analyse vectors for motion and geometric relationships.", "Show two arrows on a grid and identify components and magnitude.", "Convert to components, practise operations, and apply to motion or forces.", "Confusing magnitude and direction."],
        "general": [f"{topic}, key concepts, justification, problem solving, mathematical communication", f"Build deep understanding of {topic} through guided teaching and progressive practice.", f"Retrieval question connected to prerequisites for {topic}.", f"Worked example, guided practice, and higher-order task on {topic}.", "A common misconception identified from starter responses."],
    }
    k, focus, starter, main, misconception = (ar_bank if lang == "ar" else en_bank)[fam]
    if sk:
        k = f"{k}، {sk}" if lang == "ar" else f"{k}, {sk}"

    if lang == "ar":
        note_line = f"\nملاحظة المعلم: {note}" if note else ""
        source_line = f"\nتم استخدام النص المرفق لاختيار أمثلة وأنشطة مناسبة." if lesson.source_text else ""
        file_line = f"\nمرفق: {lesson.source_file_name}." if lesson.source_file_name and not lesson.source_text else ""
        return {
            "subject": subject, "class_name": class_name,
            "keywords": k,
            "sdg": "SDG 4 التعليم الجيد + SDG 11 مدن ومجتمعات مستدامة: توظيف التفكير الرياضي لاتخاذ قرارات دقيقة ومسؤولة.",
            "strategies": "تعلم نشط، نمذجة تفكير المعلم، Think-Pair-Share، أسئلة متدرجة، ألواح صغيرة، تغذية راجعة فورية، وتحدي للموهوبين." + note_line + source_line + file_line,
            "intervention": f"دعم فوري عبر مثال جزئي، بطاقة خطوات، شريك داعم، وأسئلة قصيرة متدرجة. خطأ متوقع: {misconception} إذا أخفق أكثر من 25% في AFL ينفذ المعلم إعادة تدريس لمدة 5 دقائق.",
            "learning_outcomes": f"بنهاية الدرس سيكون الطلاب قادرين على:\n1. يحدد المفهوم الرئيس في {topic} بدقة.\n2. يطبق القاعدة/الإجراء المناسب لحل مسائل مباشرة.\n3. يفسر الناتج رياضياً أو بيانياً.\n4. يبرر الحل بلغة رياضية صحيحة.\n5. يحل مسألة سياقية مرتبطة بالدرس.\n6. يقيم معقولية الإجابة ويصحح خطأً شائعاً.",
            "differentiation": "دعم: بطاقات خطوات وأمثلة محلولة.\nمستوى متوسط: تدريب موجه بتدرج.\nمتقدمون: سؤال HOTS يتطلب تعميماً أو تبريراً.\nIEP/APL: تبسيط اللغة وتقليل الحمل الحسابي عند الحاجة.",
            "success_criteria": f"أستطيع أن: 1. أشرح فكرة {topic}. 2. أختار الاستراتيجية المناسبة. 3. أحل بخطوات منظمة. 4. أفسر معنى الحل. 5. أصحح خطأً شائعاً. 6. أحقق 80% فأكثر في Exit Ticket.",
            "starter": starter,
            "main": main,
            "teacher_led": f"يقدم المعلم نموذجاً محلولاً لدرس {topic}، يبرز الكلمات المفتاحية، ويستخدم أسئلة تحقق قصيرة بعد كل خطوة لجمع دليل تعلم فوري.",
            "student_led": "عمل فردي ثم زوجي: يحل الطلاب مسائل متدرجة، يقارنون الطرق، ويكتبون تبريراً مختصراً. طالب قائد يشرح خطوة واحدة للمجموعة.",
            "plenary": "Exit Ticket: سؤال مهاري + سؤال تفسير + تصحيح خطأ شائع. مشاركة إجابتين مميزتين ومعالجة خطأ واحد قبل إنهاء الحصة.",
            "kpi": f"مهمة AFL على Classroom Monitor: 4 أسئلة قصيرة حول {topic}: مهاري، تمثيل/تفسير، خطأ شائع، وسؤال تحدي. تحليل النتائج لتحديد التدخل القادم.",
            "resources": "سبورة ذكية، جهاز عرض، Desmos/GeoGebra عند الحاجة، ورقة عمل قصيرة، بطاقات خطوات، آلة حاسبة، ألواح صغيرة، Classroom Monitor.",
            "identity": "ربط بالهوية الإماراتية: أمثلة من البنية التحتية أو الطاقة أو النقل أو الاستدامة في دولة الإمارات لإظهار قيمة الرياضيات في التطوير الوطني.",
            "competency": "☒ Collaboration/Teamwork   ☒ Critical Thinking   ☒ Problem Solving\n☒ Communication   ☒ Digital Competence   ☒ Innovation\n☒ Initiative/Self-direction   ☒ Global/environmental awareness",
            "curriculum": f"المادة: {subject} | الصف: {class_name}\nروابط قبلية: مهارات الجبر والتمثيل البياني.\nروابط لاحقة: تطبيق {topic} في مسائل مركبة وتقييمية.",
        }
    return {
        "subject": subject, "class_name": class_name,
        "keywords": k,
        "sdg": "SDG 4 Quality Education + SDG 11 Sustainable Cities: using mathematical thinking to make accurate and responsible decisions.",
        "strategies": "Active learning, teacher think-aloud modelling, Think-Pair-Share, tiered questioning, mini-whiteboards, immediate feedback, and gifted challenge." + (f"\nTeacher note: {note}" if note else "") + ("\nUploaded text was used to select examples and activities." if lesson.source_text else ""),
        "intervention": f"Immediate support through a partially worked example, step card, supportive partner, and short tiered questions. Anticipated misconception: {misconception} If over 25% miss the AFL check, the teacher delivers a 5-minute re-teach.",
        "learning_outcomes": f"By the end of the lesson, students will be able to:\n1. Identify the core concept in {topic} accurately.\n2. Apply a suitable rule or strategy to direct problems.\n3. Interpret the result mathematically or graphically.\n4. Justify the solution using accurate mathematical language.\n5. Solve a contextual problem linked to the lesson.\n6. Evaluate reasonableness and correct a common error.",
        "differentiation": "Support: step cards and worked examples.\nCore: guided practice with progressive difficulty.\nAdvanced: HOTS question requiring generalisation or justification.\nIEP/APL: simplified language and reduced computational load when required.",
        "success_criteria": f"I can: 1. Explain {topic}. 2. Choose the correct strategy. 3. Solve with organised steps. 4. Interpret the solution. 5. Correct a common error. 6. Achieve 80% or more in the Exit Ticket.",
        "starter": starter,
        "main": main,
        "teacher_led": f"The teacher models a worked example for {topic}, highlights key vocabulary, and uses quick checks after each step to collect learning evidence.",
        "student_led": "Individual then paired work: students solve tiered questions, compare methods, and write a short justification. A student leader explains one step to the group.",
        "plenary": "Exit Ticket: one procedural question, one reasoning question, and one error-correction prompt. Share two strong responses and address one misconception.",
        "kpi": f"Classroom Monitor AFL task: 4 short questions on {topic}: skill, representation/interpretation, common error, and challenge. Results inform the next intervention.",
        "resources": "Interactive board, projector, Desmos/GeoGebra if needed, short worksheet, step cards, calculator, mini-whiteboards, Classroom Monitor.",
        "identity": "UAE identity link: examples from infrastructure, energy, transport, or sustainability in the UAE to show mathematics in national development.",
        "competency": "☒ Collaboration/Teamwork   ☒ Critical Thinking   ☒ Problem Solving\n☒ Communication   ☒ Digital Competence   ☒ Innovation\n☒ Initiative/Self-direction   ☒ Global/environmental awareness",
        "curriculum": f"Subject: {subject} | Class: {class_name}\nPrerequisite links: algebra and graph interpretation.\nNext links: applying {topic} in compound assessment-style problems.",
    }


def ai_system_instructions(lang: str) -> str:
    if lang == "ar":
        return (
            "أنت خبير تحضير دروس لمدارس الإمارات الخاصة. اكتب خطة درس دقيقة ومختصرة وجاهزة لوضعها داخل قالب Word رسمي. "
            "اللغة عربية احترافية RTL، مع استخدام English keywords عند الحاجة فقط. "
            "المطلوب JSON فقط بدون Markdown. اجعل كل الحقول موجزة حتى لا تخرج من خلايا الجدول. "
            "يجب أن تكون نواتج التعلم ومعايير النجاح متغيرة فعلاً حسب عنوان الدرس، وليست عامة. "
            "أضف UAE identity و Sustainability و STEM و Differentiation و AFL بدقة."
        )
    return (
        "You are an expert lesson-planning assistant for UAE private schools. Generate concise, table-ready lesson plan content for an official Word template. "
        "Return JSON only with no markdown. All fields must be lesson-specific, measurable, and concise enough to fit Word table cells. "
        "Include UAE identity, sustainability, STEM, differentiation, and AFL evidence."
    )


def ai_prompt(lesson: LessonInput) -> str:
    lang_name = "Arabic" if lesson.language == "ar" else "English"
    schema_hint = {k: f"string, max {MAX_FIELD_CHARS[k]} chars" for k in REQUIRED_FIELDS}
    return json.dumps({
        "task": "Generate official school lesson plan content as JSON only.",
        "language": lang_name,
        "teacher": lesson.teacher,
        "subject": lesson.subject or ("Mathematics" if lesson.language == "en" else "رياضيات"),
        "class": lesson.class_name,
        "topic": lesson.topic,
        "date": fmt_date(lesson.date),
        "periods": lesson.periods,
        "teacher_notes": lesson.notes,
        "uploaded_lesson_text": clean_text(lesson.source_text, 2500),
        "required_json_keys": REQUIRED_FIELDS,
        "field_length_limits": schema_hint,
        "quality_rules": [
            "Learning outcomes must have exactly 6 measurable outcomes.",
            "Success criteria must have exactly 6 student-friendly criteria.",
            "Starter, main, teacher-led, student-led, and plenary must be specific to the lesson topic.",
            "Do not include generic filler text.",
            "Do not mention that you are an AI.",
            "Keep text concise for Word table cells."
        ]
    }, ensure_ascii=False)


def extract_json_from_text(text: str) -> Dict[str, Any]:
    text = clean_text(text)
    # Remove code fences if the model used them accidentally.
    text = re.sub(r"^```(?:json)?", "", text.strip(), flags=re.I)
    text = re.sub(r"```$", "", text.strip())
    try:
        return json.loads(text)
    except Exception:
        pass
    start = text.find("{")
    end = text.rfind("}")
    if start >= 0 and end > start:
        return json.loads(text[start : end + 1])
    raise ValueError("AI response did not contain valid JSON.")


def ai_content(lesson: LessonInput) -> Tuple[Dict[str, str], str]:
    # Returns (content, mode_label)
    if not os.getenv("OPENAI_API_KEY") or OpenAI is None:
        return offline_content(lesson), "offline"
    try:
        client = OpenAI()
        model = os.getenv("OPENAI_MODEL", "gpt-5.5")
        response = client.responses.create(
            model=model,
            instructions=ai_system_instructions(lesson.language),
            input=ai_prompt(lesson),
            max_output_tokens=2600,
            store=False,
        )
        output = getattr(response, "output_text", "") or ""
        data = extract_json_from_text(output)
        fallback = offline_content(lesson)
        normalized: Dict[str, str] = {
            "subject": lesson.subject or fallback.get("subject", ""),
            "class_name": lesson.class_name or fallback.get("class_name", ""),
        }
        for key in REQUIRED_FIELDS:
            value = data.get(key) or fallback.get(key, "")
            normalized[key] = clean_text(value, MAX_FIELD_CHARS.get(key, 600))
        return normalized, "ai"
    except Exception as exc:
        logger.exception("AI generation failed; falling back to offline mode")
        content = offline_content(lesson)
        if lesson.language == "ar":
            content["strategies"] += "\nتنبيه: تعذر الاتصال بالذكاء الاصطناعي، فتم استخدام الوضع الداخلي الذكي."
        else:
            content["strategies"] += "\nNote: AI connection failed, so the smart internal mode was used."
        return content, "offline_after_ai_error"


# Cache content per request to avoid double generation between preview and doc export is intentionally not used.
# Each generate call creates a fresh plan.

def build_content(lesson: LessonInput) -> Dict[str, str]:
    content, mode = ai_content(lesson)
    content["_mode"] = mode
    return content


# ----------------------------- DOCX generation -----------------------------

def generate_docx(lesson: LessonInput) -> bytes:
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError("Template file not found: assets/Lesson_Plan_Template_AY2026_2027.docx")
    lang = lesson.language
    content = build_content(lesson)
    doc = Document(str(TEMPLATE_PATH))
    date_text = fmt_date(lesson.date)
    topic = lesson.topic.strip()
    teacher = lesson.teacher.strip()
    periods = lesson.periods.strip() or ("حصة واحدة (45 دقيقة)" if lang == "ar" else "1 period (45 min)")

    t0 = doc.tables[0]
    set_cell_text(t0.cell(1, 1), teacher, lang, size=9.2)
    set_cell_text(t0.cell(1, 4), content.get("subject", lesson.subject), lang, size=9.2)
    set_cell_text(t0.cell(1, 6), date_text, lang, size=9.2)
    set_cell_text(t0.cell(2, 1), content.get("class_name", lesson.class_name), lang, size=9.2)
    set_cell_text(t0.cell(2, 4), topic, lang, size=9.2)
    set_cell_text(t0.cell(2, 6), periods, lang, size=9.2)

    set_cell_text(t0.cell(4, 2), content["keywords"], lang, size=8.0)
    set_cell_text(t0.cell(5, 2), content["sdg"], lang, size=7.7)
    set_cell_text(t0.cell(6, 2), content["strategies"], lang, size=7.0)
    set_cell_text(t0.cell(6, 5), content["intervention"], lang, size=6.8)
    set_cell_text(t0.cell(7, 2), content["learning_outcomes"], lang, size=7.0)
    set_cell_text(t0.cell(7, 5), content["differentiation"], lang, size=6.8)
    set_cell_text(t0.cell(8, 2), content["success_criteria"], lang, size=6.8)

    set_cell_text(t0.cell(10, 2), content["starter"], lang, size=7.35)
    set_cell_text(t0.cell(11, 2), content["main"], lang, size=7.1)
    set_cell_text(t0.cell(12, 2), content["teacher_led"], lang, size=7.1)
    set_cell_text(t0.cell(13, 2), content["student_led"], lang, size=7.1)
    set_cell_text(t0.cell(14, 2), content["plenary"], lang, size=7.1)

    t1 = doc.tables[1]
    set_cell_text(t1.cell(0, 1), content["kpi"], lang, size=7.0)
    set_cell_text(t1.cell(0, 2), f"Date of Upload:\n{date_text}" if lang == "en" else f"تاريخ الرفع:\n{date_text}", lang, size=7.0)
    set_cell_text(t1.cell(0, 3), "Time of Upload:\nBefore next lesson" if lang == "en" else "وقت الرفع:\nقبل الحصة التالية", lang, size=7.0)
    set_cell_text(t1.cell(0, 5), "IEP/APL: attached if required" if lang == "en" else "IEP/APL: يرفق عند الحاجة", lang, size=7.0)
    set_cell_text(t1.cell(1, 1), "Creativity (✓)" if lang == "en" else "الإبداع (✓)", lang, size=7.7)
    set_cell_text(t1.cell(1, 2), "Collaboration (✓)" if lang == "en" else "التعاون (✓)", lang, size=7.7)
    set_cell_text(t1.cell(1, 3), "Critical Thinking (✓)" if lang == "en" else "التفكير الناقد (✓)", lang, size=7.7)
    set_cell_text(t1.cell(1, 5), "Communication (✓)" if lang == "en" else "التواصل (✓)", lang, size=7.7)
    set_cell_text(t1.cell(2, 1), content["resources"], lang, size=7.0)
    # In the official template this cell is the curriculum link box on page 3.
    try:
        set_cell_text(t1.cell(3, 4), content["curriculum"], lang, size=6.8)
    except Exception:
        pass

    t2 = doc.tables[2]
    for row_idx in range(1, 8):
        fill_check(t2.cell(row_idx, 1), lang=lang)
    for row_idx in range(1, 5):
        fill_check(t2.cell(row_idx, 3), lang=lang)
    for row_idx in range(1, 8):
        fill_check(t2.cell(row_idx, 5), lang=lang)
    fill_check(t2.cell(7, 5), lang=lang)

    # Metadata: remove old app name and keep it clean.
    doc.core_properties.title = f"Lesson Plan - {topic}"
    doc.core_properties.subject = content.get("subject", lesson.subject)
    doc.core_properties.author = teacher or "Magdy Lesson Planner"
    doc.core_properties.keywords = content.get("keywords", "")
    doc.core_properties.comments = "Generated by Magdy Lesson Planner"

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# ----------------------------- Request parsing -----------------------------

def parse_lessons_from_request() -> Tuple[List[LessonInput], List[str]]:
    errors: List[str] = []
    teacher = request.form.get("teacher", "").strip()
    subject = request.form.get("subject", "").strip()
    class_name = request.form.get("class_name", "").strip()
    periods = request.form.get("periods", "").strip()
    language = request.form.get("language", "ar")
    access_code = request.form.get("access_code", "").strip()
    pasted_image = request.form.get("pasted_image", "")

    ok_code, code_message = validate_access_code(access_code)
    if not ok_code:
        errors.append(code_message)

    topics = request.form.getlist("topic[]") or [request.form.get("topic", "")]
    dates = request.form.getlist("date[]") or [request.form.get("date", "")]
    notes_list = request.form.getlist("lesson_notes[]") or [request.form.get("lesson_notes", "")]
    source_text_list = request.form.getlist("source_text[]") or [request.form.get("source_text", "")]

    max_len = max(len(topics), len(dates), len(notes_list), len(source_text_list))
    lessons: List[LessonInput] = []
    for idx in range(max_len):
        topic = topics[idx].strip() if idx < len(topics) else ""
        date = dates[idx].strip() if idx < len(dates) else ""
        notes = notes_list[idx].strip() if idx < len(notes_list) else ""
        source_text = source_text_list[idx].strip() if idx < len(source_text_list) else ""
        if not topic and not date and not notes and not source_text:
            continue
        if not topic:
            errors.append(f"Lesson {idx+1}: title is required.")

        fname = ""
        file_key = f"lesson_file_{idx}"
        uploaded_name, extracted = save_uploaded_file(request.files.get(file_key), f"lesson{idx+1}")
        if uploaded_name:
            fname = uploaded_name
        if extracted and not extracted.startswith("[Unsupported"):
            source_text = (source_text + "\n" + extracted).strip()
        if extracted.startswith("[Unsupported"):
            errors.append(extracted)
        if idx == 0 and pasted_image:
            pasted_name = save_pasted_image(pasted_image, "pasted")
            if pasted_name and not fname:
                fname = pasted_name

        lessons.append(LessonInput(
            index=idx + 1,
            teacher=teacher,
            subject=subject,
            class_name=class_name,
            periods=periods,
            language="ar" if language == "ar" else "en",
            topic=topic,
            date=date,
            notes=notes,
            source_text=source_text,
            source_file_name=fname,
            ai_mode="server",
            access_code=access_code,
        ))
    if not teacher:
        errors.append("Teacher name is required.")
    if not lessons:
        errors.append("Add at least one lesson title.")
    return lessons, errors


def store_docx_file(lesson: LessonInput, docx_bytes: bytes) -> str:
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    teacher = safe_filename(lesson.teacher, "teacher")
    topic = safe_filename(lesson.topic, "lesson")
    filename = f"{stamp}_Lesson_Plan_{teacher}_{topic}.docx"
    (LIBRARY_DIR / filename).write_bytes(docx_bytes)
    add_to_library({
        "filename": filename,
        "topic": lesson.topic,
        "teacher": lesson.teacher,
        "subject": lesson.subject or ("رياضيات" if lesson.language == "ar" else "Mathematics"),
        "class_name": lesson.class_name or ("الثاني عشر متقدم" if lesson.language == "ar" else "Grade 12 Advanced"),
        "date": fmt_date(lesson.date),
        "language": lesson.language,
        "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "source_file": lesson.source_file_name,
    })
    return filename



# ----------------------------- Curriculum planner helpers -----------------------------

CURRICULUM_ALLOWED_EXTENSIONS = {"pdf", "docx", "txt", "md", "csv", "pptx", "xlsx", "png", "jpg", "jpeg", "webp"}


def _curriculum_allowed(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in CURRICULUM_ALLOWED_EXTENSIONS


def _curriculum_job_path(job_id: str) -> Path:
    safe_id = re.sub(r"[^a-fA-F0-9]", "", job_id or "")[:64]
    if not safe_id:
        raise ValueError("Invalid plan session")
    return CURRICULUM_JOB_DIR / f"{safe_id}.json"


def _curriculum_save_job(job_id: str, payload: dict) -> None:
    _curriculum_job_path(job_id).write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def _curriculum_load_job(job_id: str) -> dict:
    path = _curriculum_job_path(job_id)
    if not path.exists():
        raise FileNotFoundError("انتهت جلسة الخطة. أنشئ الخطة مرة أخرى.")
    return json.loads(path.read_text(encoding="utf-8"))


def _curriculum_clean_old_files(max_age_hours: int = 24) -> None:
    cutoff = time.time() - max_age_hours * 3600
    for folder in (CURRICULUM_UPLOAD_DIR, CURRICULUM_JOB_DIR, CURRICULUM_EXPORT_DIR):
        for path in folder.glob("*"):
            try:
                if path.stat().st_mtime < cutoff:
                    path.unlink() if path.is_file() else shutil.rmtree(path)
            except OSError:
                pass


def _curriculum_apply_edits(payload: dict) -> dict:
    plan_type = payload["plan_type"]
    meta = payload["meta"]
    for key in ("teacher", "subject", "grade", "academic_year"):
        value = request.form.get(key)
        if value is not None:
            meta[key] = value.strip()
    data = payload["plan"]
    if plan_type == "medium":
        for idx, week in enumerate(data["weeks"]):
            for field in ("content", "learning_objectives", "ai_literacy", "resources"):
                value = request.form.get(f"week_{idx}_{field}")
                if value is not None:
                    week[field] = value.strip()
        for field in (
            "title", "targets", "assessment_opportunities", "century_skills", "vocabulary",
            "eps_guiding_statement", "global_citizenship", "cross_curricular", "national_identity",
            "ai_integration_approach", "guardrails_prompt_controls", "cognitive_integrity_strategy", "ai_safeguarding",
        ):
            value = request.form.get(field)
            if value is not None:
                data[field] = value.strip()
    else:
        for idx, half in enumerate(data["half_terms"]):
            for field in ("content", "summative_assessment"):
                value = request.form.get(f"half_{idx}_{field}")
                if value is not None:
                    half[field] = value.strip()
    payload["meta"] = meta
    payload["plan"] = data
    return payload


# ----------------------------- Routes -----------------------------


@app.errorhandler(413)
def too_large(_error):
    message = "حجم الملف أكبر من 35 MB."
    if request.path.startswith("/curriculum"):
        flash(message, "error")
        return redirect(url_for("curriculum_planner"))
    return render_template("index.html", error=message, status=status_payload()), 413


@app.errorhandler(Exception)
def handle_exception(exc):
    logger.exception("Unhandled application error")
    message = "حدث خطأ داخل التطبيق ولم يتم إغلاق السيرفر. راجع ملف logs/error_log.txt."
    if request.path.startswith("/api/"):
        return jsonify({"ok": False, "error": message, "details": str(exc)}), 500
    if request.path.startswith("/curriculum"):
        return render_template("curriculum_planner.html", error=f"{message}\n{exc}", status=status_payload(), initial_plan_type="medium"), 500
    return render_template("index.html", error=f"{message}\n{exc}", status=status_payload()), 500



@app.route("/")
def dashboard():
    _curriculum_clean_old_files()
    return render_template("dashboard.html", status=status_payload())


@app.route("/lesson-planner")
def index():
    return render_template("index.html", status=status_payload())


@app.get("/curriculum-planner")
def curriculum_planner():
    _curriculum_clean_old_files()
    initial_plan_type = request.args.get("type", "medium")
    if initial_plan_type not in {"medium", "long"}:
        initial_plan_type = "medium"
    return render_template("curriculum_planner.html", status=status_payload(), initial_plan_type=initial_plan_type)


@app.post("/curriculum/generate")
def curriculum_generate():
    started_at = time.monotonic()
    plan_type = request.form.get("plan_type", "medium")
    raw_language = request.form.get("language", "Arabic")
    language = "Arabic" if str(raw_language).strip().casefold() in {"arabic", "ar", "العربية", "عربي"} else "English"
    meta = {
        "teacher": request.form.get("teacher", "").strip(),
        "subject": request.form.get("subject", "").strip(),
        "grade": request.form.get("grade", "").strip(),
        "academic_year": request.form.get("academic_year", "2026-2027").strip(),
        "language": language,
    }
    if not all([meta["teacher"], meta["subject"], meta["grade"]]):
        flash("يرجى إدخال اسم المعلم والمادة والصف.", "error")
        return redirect(url_for("curriculum_planner", type=plan_type))

    manual_topics = request.form.get("manual_topics", "").strip()
    instructions = request.form.get("instructions", "").strip()
    uploaded = request.files.get("curriculum_file")
    source_text = ""
    source_name = "Manual topics"
    if uploaded and uploaded.filename:
        if not _curriculum_allowed(uploaded.filename):
            flash("نوع الملف غير مدعوم. استخدم PDF أو Word أو PowerPoint أو Excel أو صورة أو ملف نصي.", "error")
            return redirect(url_for("curriculum_planner", type=plan_type))
        safe_upload_name = secure_filename(uploaded.filename) or f"curriculum.{file_ext(uploaded.filename) or 'txt'}"
        filename = f"{uuid.uuid4().hex}_{safe_upload_name}"
        upload_path = CURRICULUM_UPLOAD_DIR / filename
        uploaded.save(upload_path)
        source_name = uploaded.filename
        try:
            source_text = extract_curriculum_text(upload_path)
        except Exception as exc:
            logger.exception("Curriculum file extraction failed")
            flash(f"تعذر قراءة الملف: {exc}", "error")
            return redirect(url_for("curriculum_planner", type=plan_type))

    if not source_text and not manual_topics:
        flash("ارفع الكتاب أو الفهرس، أو الصق قائمة الموضوعات.", "error")
        return redirect(url_for("curriculum_planner", type=plan_type))

    topics = candidate_topics(source_text, manual_topics)
    logger.info("Curriculum extraction complete: type=%s source=%s chars=%s candidates=%s elapsed=%.2fs", plan_type, source_name, len(source_text), len(topics), time.monotonic() - started_at)
    topics = refine_topics(meta, source_text, topics, language)
    logger.info("Curriculum topics ready: type=%s topics=%s elapsed=%.2fs", plan_type, len(topics), time.monotonic() - started_at)
    if not topics:
        flash("لم يتم اكتشاف عناوين وحدات أو دروس واضحة. ارفع صفحة الفهرس أو الصق الموضوعات يدويًا.", "error")
        return redirect(url_for("curriculum_planner", type=plan_type))

    try:
        if plan_type == "long":
            plan = generate_long(meta, source_text, topics, language, instructions)
        else:
            plan_type = "medium"
            plan = generate_medium(meta, source_text, topics, language, instructions)
    except Exception as exc:
        logger.exception("Curriculum plan generation failed")
        flash(f"تعذر إنشاء الخطة: {exc}", "error")
        return redirect(url_for("curriculum_planner", type=plan_type))

    job_id = uuid.uuid4().hex
    payload = {
        "job_id": job_id,
        "plan_type": plan_type,
        "meta": meta,
        "plan": plan.model_dump(),
        "source_name": source_name,
        "detected_topics": topics[:50],
    }
    _curriculum_save_job(job_id, payload)
    logger.info("Curriculum plan ready: type=%s job=%s elapsed=%.2fs", plan_type, job_id[:8], time.monotonic() - started_at)
    return render_template("curriculum_preview.html", job=payload, status=status_payload())


@app.post("/curriculum/export/<job_id>/<fmt>")
def curriculum_export(job_id: str, fmt: str):
    try:
        payload = _curriculum_apply_edits(_curriculum_load_job(job_id))
        _curriculum_save_job(job_id, payload)
        meta = payload["meta"]
        safe_subject = safe_filename(meta["subject"], "Subject")
        if payload["plan_type"] == "medium":
            plan = MediumPlan.model_validate(payload["plan"])
            docx_path = CURRICULUM_EXPORT_DIR / f"MTP_{safe_subject}_{job_id[:8]}.docx"
            build_medium(meta, plan, docx_path)
        else:
            plan = LongPlan.model_validate(payload["plan"])
            docx_path = CURRICULUM_EXPORT_DIR / f"LTP_{safe_subject}_{job_id[:8]}.docx"
            build_long(meta, plan, docx_path)

        if fmt == "pdf":
            pdf_path = CURRICULUM_EXPORT_DIR / f"{docx_path.stem}.pdf"
            if payload["plan_type"] == "medium":
                build_medium_pdf(meta, plan, pdf_path)
            else:
                build_long_pdf(meta, plan, pdf_path)
            return send_file(pdf_path, as_attachment=True, download_name=pdf_path.name, mimetype="application/pdf")
        return send_file(
            docx_path,
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            as_attachment=True,
            download_name=docx_path.name,
        )
    except Exception as exc:
        logger.exception("Curriculum export failed")
        flash(f"تعذر تصدير الخطة: {exc}", "error")
        return redirect(url_for("curriculum_planner"))


@app.route("/api/status")
def api_status():
    return jsonify(status_payload())


@app.route("/api/preview", methods=["POST"])
def preview():
    lessons, errors = parse_lessons_from_request()
    if errors:
        return jsonify({"ok": False, "errors": errors}), 400
    lesson = lessons[0]
    content = build_content(lesson)
    return jsonify({"ok": True, "lesson": asdict(lesson), "content": content, "status": status_payload()})


@app.route("/generate", methods=["POST"])
def generate():
    lessons, errors = parse_lessons_from_request()
    if errors:
        return render_template("index.html", error=" | ".join(errors), status=status_payload()), 400

    user_key = lessons[0].teacher or request.remote_addr or "anonymous"
    for _lesson in lessons:
        ok_limit, limit_msg = check_usage_limit(user_key)
        if not ok_limit:
            return render_template("index.html", error=limit_msg, status=status_payload()), 429

    generated: List[Tuple[str, bytes]] = []
    try:
        for lesson in lessons:
            docx_bytes = generate_docx(lesson)
            stored_name = store_docx_file(lesson, docx_bytes)
            generated.append((stored_name, docx_bytes))
    except Exception as exc:
        logger.exception("Error while generating DOCX")
        return render_template(
            "index.html",
            error=f"حدث خطأ أثناء إنشاء الملف، لكن التطبيق لم يغلق. التفاصيل محفوظة في logs/error_log.txt\n{exc}",
            status=status_payload(),
        ), 500

    if len(generated) == 1:
        filename, docx_bytes = generated[0]
        return send_file(
            io.BytesIO(docx_bytes),
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            as_attachment=True,
            download_name=filename,
        )

    zip_bytes = io.BytesIO()
    with zipfile.ZipFile(zip_bytes, "w", zipfile.ZIP_DEFLATED) as zf:
        for filename, docx_bytes in generated:
            zf.writestr(filename, docx_bytes)
    zip_bytes.seek(0)
    zip_name = f"Lesson_Plans_Batch_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
    return send_file(zip_bytes, mimetype="application/zip", as_attachment=True, download_name=zip_name)


@app.route("/library")
def library():
    items = [x for x in load_library() if (LIBRARY_DIR / x.get("filename", "")).exists()]
    return render_template("library.html", items=items, status=status_payload())


@app.route("/download/<path:filename>")
def download(filename: str):
    safe = os.path.basename(filename)
    path = LIBRARY_DIR / safe
    if not path.exists():
        return "File not found", 404
    return send_file(str(path), as_attachment=True, download_name=safe)


@app.route("/delete/<path:filename>", methods=["POST"])
def delete(filename: str):
    safe = os.path.basename(filename)
    path = LIBRARY_DIR / safe
    if path.exists():
        path.unlink()
    items = [x for x in load_library() if x.get("filename") != safe]
    save_library(items)
    return redirect(url_for("library"))


@app.route("/health")
def health():
    return {"status": "ok", "time": time.time(), **status_payload()}


if __name__ == "__main__":
    host = os.getenv("HOST", "127.0.0.1")
    port = env_int("PORT", 5000)
    app.run(host=host, port=port, debug=False, threaded=True)
