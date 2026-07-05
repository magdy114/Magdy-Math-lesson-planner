from __future__ import annotations

import os
import re
import time
from pathlib import Path

TEXT_LIMIT = 28_000
TOPIC_LIMIT = 90

NOISE_TERMS = (
    "copyright", "all rights reserved", "rights reserved", "permission", "publisher",
    "publishing", "isbn", "mcgraw", "pearson", "education, llc", "trademark",
    "sourced from", "source:", "printed in", "edition", "www.", "http://", "https://",
    "©", "®", "™", "differentiation", "teacher edition", "student edition",
    "حقوق الطبع", "جميع الحقوق محفوظة", "الناشر", "الطبعة", "ردمك", "حقوق النشر",
    "تم النشر", "لا يجوز", "إعادة إنتاج", "المؤلف", "المراجع", "المصدر",
)

HEADING_WORDS = (
    "unit", "chapter", "lesson", "section", "module", "topic", "strand", "domain",
    "الوحدة", "الفصل", "الدرس", "الموضوع", "المحور", "المجال", "الباب", "القسم",
)

BODY_HINTS = (
    "example", "worked example", "exercise", "practice", "find the", "calculate", "suppose",
    "if the", "where the", "is defined", "learning objective", "success criteria", "warning",
    "note that", "notice", "figure", "table", "answer", "solution", "check your understanding",
    "حل المثال", "مثال", "تمرين", "تدرب", "أوجد", "احسب", "إذا كان", "حيث إن",
    "يوضح الشكل", "الشكل", "الجدول", "نشاط", "ناقش", "فسر", "علل", "اختر الإجابة",
    "اكتب", "استخدم الشكل", "الحل", "تحذير", "لاحظ", "تذكر", "بما أن", "نفترض",
)

GENERIC_NON_TOPICS = {
    "contents", "table of contents", "index", "glossary", "references", "acknowledgements",
    "introduction", "preface", "answers", "answer key", "review", "assessment",
    "guided practice", "independent practice", "homework", "objectives", "resources",
    "الفهرس", "المحتويات", "المقدمة", "المراجع", "الإجابات", "مفتاح الإجابة",
    "مراجعة", "تقويم", "أهداف التعلم", "المصادر", "الواجب", "التدريب الموجه",
}

ARABIC_DIGITS = str.maketrans("٠١٢٣٤٥٦٧٨٩", "0123456789")


def _env_int(name: str, default: int, minimum: int, maximum: int) -> int:
    try:
        value = int(os.getenv(name, str(default)))
    except Exception:
        value = default
    return max(minimum, min(maximum, value))


def _env_float(name: str, default: float, minimum: float, maximum: float) -> float:
    try:
        value = float(os.getenv(name, str(default)))
    except Exception:
        value = default
    return max(minimum, min(maximum, value))


def _normalize(text: str) -> str:
    text = (text or "").replace("\x00", " ")
    text = re.sub(r"[\t\u00a0]+", " ", text)
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ ]{2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()[:TEXT_LIMIT]


def _remove_numbering(text: str) -> str:
    """Remove lesson/page numbering without damaging the teachable title."""
    text = text.translate(ARABIC_DIGITS)
    text = re.sub(r"^\s*\[HEADING\]\s*", "", text, flags=re.I)
    # Labels and codes at the beginning: Lesson 1-3, الدرس 3.2, Chapter 4 ...
    labels = r"unit|chapter|lesson|section|module|topic|الوحدة|الفصل|الدرس|القسم|الباب|الموضوع"
    text = re.sub(
        rf"^\s*(?:{labels})\s*(?:no\.?|رقم)?\s*[0-9]+(?:[.\-][0-9]+)*\s*[:\-–—|]*\s*",
        "", text, flags=re.I,
    )
    # Hierarchical code at the beginning: 3.2 Title / 1-3 Title.
    text = re.sub(r"^\s*[0-9]+(?:[.\-][0-9]+){1,3}\s*[:\-–—|]*\s*", "", text)
    # Trailing lesson/page code.
    text = re.sub(
        rf"\s*[|\-–—:]\s*(?:{labels})?\s*[0-9]+(?:[.\-][0-9]+)*\s*$",
        "", text, flags=re.I,
    )
    text = re.sub(r"\.{2,}\s*[0-9]+\s*$", "", text)
    text = re.sub(r"\s+(?:page|p\.?|صفحة)\s*[0-9]+\s*$", "", text, flags=re.I)
    # Remove isolated numeric tokens. A professional plan displays titles, not book numbering.
    if os.getenv("CURRICULUM_KEEP_TOPIC_NUMBERS", "0").strip().lower() not in {"1", "true", "yes"}:
        text = re.sub(r"(?<![A-Za-z])\b[0-9]+(?:[.\-][0-9]+)*\b(?![A-Za-z])", " ", text)
    return re.sub(r"\s+", " ", text).strip(" -–—|:;,.")


def _clean_title(text: str) -> str:
    text = _remove_numbering(str(text or ""))
    text = re.sub(r"^[•▪◦●✓✔]+\s*", "", text)
    text = re.sub(r"\s+", " ", text).strip(" -–—|:;,.")
    return text


def _is_noise(text: str) -> bool:
    clean = _clean_title(text)
    lower = clean.casefold()
    if not clean or lower in GENERIC_NON_TOPICS:
        return True
    if any(term in lower for term in NOISE_TERMS):
        return True
    if "@" in clean or re.search(r"\.(?:com|org|net|edu|ae)\b", lower):
        return True
    if re.search(r"\b(?:ISBN|DOI)\b", clean, re.I):
        return True
    if re.fullmatch(r"[\d\W_]+", clean):
        return True
    return False


def _looks_like_body(text: str) -> bool:
    clean = _clean_title(text)
    lower = clean.casefold()
    words = clean.split()
    if any(hint in lower for hint in BODY_HINTS):
        return True
    if len(clean) > 105 or len(words) > 13:
        return True
    if len(words) >= 7 and re.search(r"[.!?؟؛:]$", clean):
        return True
    digit_count = sum(ch.isdigit() for ch in clean)
    operator_count = sum(ch in "+=<>[]{}()" for ch in clean)
    if digit_count >= 3 or operator_count >= 3:
        return True
    if re.search(r"\b(?:km|mph|m/s|cm|mm|kg|sec|seconds?|دقيقة|ثانية|متر|ميل)\b", lower) and len(words) > 5:
        return True
    return False


def _is_explicit_heading(raw: str, clean: str) -> bool:
    lower_raw = str(raw or "").lstrip().casefold()
    lower = clean.casefold()
    if lower_raw.startswith("[heading]"):
        return True
    if any(re.match(rf"^{re.escape(word)}\b", lower_raw) for word in HEADING_WORDS):
        return True
    if re.match(r"^\s*[0-9٠-٩]+(?:[.\-][0-9٠-٩]+){1,3}\s+\S", str(raw or "")):
        return True
    return False


def _title_shape_ok(clean: str) -> bool:
    words = clean.split()
    if not (1 <= len(words) <= 13):
        return False
    if len(clean) > 105 or _is_noise(clean) or _looks_like_body(clean):
        return False
    # A title normally does not end as a full sentence.
    if re.search(r"[.!?؟؛]$", clean):
        return False
    return True


def _dedupe_key(text: str) -> str:
    return re.sub(r"[^\w\u0600-\u06ff]+", "", _clean_title(text).casefold())


def _append_unique(output: list[str], seen: set[str], raw: str) -> None:
    clean = _clean_title(raw)
    if not _title_shape_ok(clean):
        return
    key = _dedupe_key(clean)
    if not key or key in seen:
        return
    # Avoid titles that are mostly contained in a longer, more descriptive title already kept.
    if any(key in old or old in key for old in seen if min(len(key), len(old)) >= 10):
        return
    seen.add(key)
    output.append(clean)


def _pdf_line_candidates(page) -> list[tuple[str, float, float, bool]]:
    """Return (text, max size, page median size, bold) with bounded processing."""
    data = page.get_text("dict", sort=True)
    raw: list[tuple[str, float, bool]] = []
    sizes: list[float] = []
    for block in data.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            spans = line.get("spans", [])
            text = "".join(str(span.get("text", "")) for span in spans).strip()
            if not text:
                continue
            line_sizes = [float(span.get("size", 0) or 0) for span in spans if str(span.get("text", "")).strip()]
            if not line_sizes:
                continue
            sizes.extend(line_sizes)
            bold = any(
                "bold" in str(span.get("font", "")).casefold() or int(span.get("flags", 0) or 0) & 16
                for span in spans
            )
            raw.append((text, max(line_sizes), bold))
    median = sorted(sizes)[len(sizes) // 2] if sizes else 10.0
    return [(text, size, median, bold) for text, size, bold in raw]


def extract_pdf(path: Path) -> str:
    """Extract only likely curriculum headings using bookmarks and visible typography."""
    import fitz

    max_pages = _env_int("CURRICULUM_PDF_MAX_PAGES", 42, 8, 70)
    deadline = time.monotonic() + _env_float("CURRICULUM_EXTRACT_SECONDS", 12.0, 5.0, 20.0)
    output: list[str] = []
    seen: set[str] = set()

    with fitz.open(path) as pdf:
        if getattr(pdf, "needs_pass", False):
            raise ValueError("ملف PDF محمي بكلمة مرور.")

        try:
            toc = pdf.get_toc(simple=True) or []
        except Exception:
            toc = []
        for item in toc[:250]:
            if len(item) >= 2:
                _append_unique(output, seen, str(item[1]))
        if len(output) >= 6:
            return _normalize("\n".join(f"[HEADING] {x}" for x in output))

        for page_no in range(min(len(pdf), max_pages)):
            if time.monotonic() >= deadline or len(output) >= TOPIC_LIMIT:
                break
            try:
                candidates = _pdf_line_candidates(pdf[page_no])
            except Exception:
                continue
            kept = 0
            for raw, max_size, median_size, bold in candidates:
                clean = _clean_title(raw)
                explicit = _is_explicit_heading(raw, clean)
                typography = (
                    _title_shape_ok(clean)
                    and (
                        max_size >= max(11.0, median_size * 1.28)
                        or (bold and max_size >= max(10.5, median_size * 1.08))
                    )
                )
                if not (explicit or typography):
                    continue
                before = len(output)
                _append_unique(output, seen, clean)
                if len(output) > before:
                    kept += 1
                if kept >= 10:
                    break

    return _normalize("\n".join(f"[HEADING] {x}" for x in output))


def extract_docx(path: Path) -> str:
    from docx import Document

    deadline = time.monotonic() + _env_float("CURRICULUM_EXTRACT_SECONDS", 12.0, 5.0, 20.0)
    doc = Document(path)
    output: list[str] = []
    seen: set[str] = set()

    for paragraph in doc.paragraphs:
        if time.monotonic() >= deadline or len(output) >= TOPIC_LIMIT:
            break
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name or "").casefold() if paragraph.style else ""
        explicit = "heading" in style or "title" in style or _is_explicit_heading(text, _clean_title(text))
        if explicit:
            _append_unique(output, seen, text)

    # Tables are commonly used for a contents page. Treat short individual cells as candidates.
    for table in doc.tables[:20]:
        if time.monotonic() >= deadline or len(output) >= TOPIC_LIMIT:
            break
        for row in table.rows[:100]:
            for cell in row.cells:
                for raw in cell.text.splitlines():
                    if _is_explicit_heading(raw, _clean_title(raw)) or _title_shape_ok(_clean_title(raw)):
                        _append_unique(output, seen, raw)
            if time.monotonic() >= deadline or len(output) >= TOPIC_LIMIT:
                break

    return _normalize("\n".join(f"[HEADING] {x}" for x in output))


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    output: list[str] = []
    seen: set[str] = set()
    deadline = time.monotonic() + _env_float("CURRICULUM_EXTRACT_SECONDS", 12.0, 5.0, 20.0)
    prs = Presentation(path)
    for slide in prs.slides[:100]:
        if time.monotonic() >= deadline or len(output) >= TOPIC_LIMIT:
            break
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None and getattr(title_shape, "text", "").strip():
            _append_unique(output, seen, title_shape.text)
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if not isinstance(text, str):
                continue
            for raw in text.splitlines():
                if _is_explicit_heading(raw, _clean_title(raw)):
                    _append_unique(output, seen, raw)
    return _normalize("\n".join(f"[HEADING] {x}" for x in output))


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    output: list[str] = []
    seen: set[str] = set()
    deadline = time.monotonic() + _env_float("CURRICULUM_EXTRACT_SECONDS", 12.0, 5.0, 20.0)
    wb = load_workbook(path, data_only=True, read_only=True)
    try:
        for ws in wb.worksheets[:12]:
            if time.monotonic() >= deadline or len(output) >= TOPIC_LIMIT:
                break
            for row_no, row in enumerate(ws.iter_rows(values_only=True), 1):
                if row_no > 350 or time.monotonic() >= deadline:
                    break
                for value in row:
                    if value is None:
                        continue
                    raw = str(value).strip()
                    if _is_explicit_heading(raw, _clean_title(raw)) or _title_shape_ok(_clean_title(raw)):
                        _append_unique(output, seen, raw)
                if len(output) >= TOPIC_LIMIT:
                    break
    finally:
        wb.close()
    return _normalize("\n".join(f"[HEADING] {x}" for x in output))


def extract_image(path: Path) -> str:
    try:
        import pytesseract
        from PIL import Image

        with Image.open(path) as image:
            image.thumbnail((2200, 2200))
            text = pytesseract.image_to_string(
                image, lang=os.getenv("TESSERACT_LANG", "eng+ara"), timeout=8
            )
        return _normalize(text)
    except Exception:
        return ""


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".xlsx":
        return extract_xlsx(path)
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return extract_image(path)
    if suffix in {".txt", ".md", ".csv"}:
        return _normalize(path.read_text(encoding="utf-8", errors="ignore"))
    raise ValueError("نوع الملف غير مدعوم")


def _manual_topic_lines(manual_topics: str) -> list[str]:
    output: list[str] = []
    seen: set[str] = set()
    for raw in re.split(r"[\n;•]+", manual_topics or ""):
        _append_unique(output, seen, raw)
        if len(output) >= TOPIC_LIMIT:
            break
    return output


def candidate_topics(source_text: str, manual_topics: str = "") -> list[str]:
    """Return a strict, ordered list of teachable titles only."""
    output = _manual_topic_lines(manual_topics)
    seen = {_dedupe_key(item) for item in output}

    for raw_line in (source_text or "").splitlines():
        fragments = [raw_line]
        if "|" in raw_line:
            fragments = [part.strip() for part in raw_line.split("|") if part.strip()]
        for raw in fragments:
            clean = _clean_title(raw)
            if _is_explicit_heading(raw, clean) or _title_shape_ok(clean):
                _append_unique(output, seen, clean)
            if len(output) >= TOPIC_LIMIT:
                return output
    return output
