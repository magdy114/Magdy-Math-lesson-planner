from __future__ import annotations

import hashlib
import json
import os
import re
import threading
import time
import uuid
from collections import Counter
from pathlib import Path
from typing import Any

from werkzeug.utils import secure_filename


SOURCE_CACHE_VERSION = "source-grounded-v5"
SOURCE_MAX_CHARS = 12_000
EXTRACT_MAX_CHARS = 28_000

_STOPWORDS = {
    "the", "and", "for", "with", "from", "that", "this", "into", "using", "use", "are", "was", "were", "will",
    "lesson", "student", "students", "teacher", "page", "unit", "chapter", "example", "practice", "activity", "learning",
    "في", "من", "على", "إلى", "الى", "عن", "هذا", "هذه", "ذلك", "تلك", "التي", "الذي", "مع", "ثم", "أو", "او",
    "درس", "الدرس", "الطلاب", "الطالب", "المعلم", "صفحة", "الوحدة", "الفصل", "مثال", "تدريب", "نشاط", "التعلم",
}

_AI_GATE = threading.BoundedSemaphore(max(1, min(4, int(os.getenv("SOURCE_AI_CONCURRENCY", "2")))))


def _clean(value: str, limit: int | None = None) -> str:
    text = str(value or "").replace("\x00", " ").replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[\t\u00a0]+", " ", text)
    text = re.sub(r"[ ]{2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text[:limit] if limit else text


def _tokens(value: str) -> list[str]:
    return [
        token.casefold()
        for token in re.findall(r"[A-Za-z][A-Za-z'-]{2,}|[\u0600-\u06ff]{3,}", str(value or ""))
        if token.casefold() not in _STOPWORDS
    ]


def _split_blocks(source_text: str) -> list[str]:
    text = _clean(source_text, EXTRACT_MAX_CHARS)
    if not text:
        return []

    raw_blocks = re.split(r"\n\s*\n|(?=\[Page\s+\d+\])|(?=Slide\s+\d+\s*:)", text, flags=re.I)
    blocks: list[str] = []
    for raw in raw_blocks:
        raw = _clean(raw)
        if not raw:
            continue
        if len(raw) <= 1400:
            blocks.append(raw)
            continue
        lines = raw.splitlines()
        current: list[str] = []
        size = 0
        for line in lines:
            if current and size + len(line) > 1200:
                blocks.append("\n".join(current))
                current, size = [], 0
            current.append(line)
            size += len(line) + 1
        if current:
            blocks.append("\n".join(current))
    return blocks


def select_relevant_source(source_text: str, topic: str, subject: str = "", notes: str = "", max_chars: int = SOURCE_MAX_CHARS) -> str:
    """Select the source passages most relevant to the lesson title without losing document order."""
    blocks = _split_blocks(source_text)
    if not blocks:
        return ""

    query_tokens = set(_tokens(f"{subject} {topic} {notes}"))
    scored: list[tuple[float, int, str]] = []
    for index, block in enumerate(blocks):
        block_tokens = _tokens(block)
        counts = Counter(block_tokens)
        overlap = sum(min(counts[token], 4) for token in query_tokens)
        phrase_bonus = 0
        topic_clean = _clean(topic).casefold()
        if topic_clean and len(topic_clean) >= 4 and topic_clean in block.casefold():
            phrase_bonus = 10
        heading_bonus = 2 if re.search(r"(?:^|\n)\s*(?:unit|chapter|lesson|section|الوحدة|الفصل|الدرس|الموضوع)\b", block, re.I) else 0
        position_bonus = max(0.0, 1.5 - index * 0.04)
        scored.append((overlap * 3 + phrase_bonus + heading_bonus + position_bonus, index, block))

    # If the title is not found, preserve the beginning and then add the strongest blocks.
    chosen: dict[int, str] = {}
    total = 0
    for index in range(min(2, len(blocks))):
        block = blocks[index]
        chosen[index] = block
        total += len(block) + 20

    for _, index, block in sorted(scored, key=lambda item: (-item[0], item[1])):
        if index in chosen:
            continue
        if total + len(block) + 20 > max_chars:
            continue
        chosen[index] = block
        total += len(block) + 20
        if total >= max_chars * 0.9:
            break

    ordered = [chosen[index] for index in sorted(chosen)]
    return _clean("\n\n".join(ordered), max_chars)


def _distinctive_terms(source_context: str, limit: int = 12) -> list[str]:
    counts = Counter(_tokens(source_context))
    return [term for term, _ in counts.most_common(limit)]


def _extract_pdf(path: Path) -> str:
    max_pages = max(4, min(40, int(os.getenv("LESSON_PDF_MAX_PAGES", "18"))))
    pieces: list[str] = []
    try:
        import fitz

        with fitz.open(path) as pdf:
            for page_index in range(min(len(pdf), max_pages)):
                page = pdf[page_index]
                text = page.get_text("text", sort=True) or ""
                if text.strip():
                    pieces.append(f"[Page {page_index + 1}]\n{text.strip()}")
            combined = _clean("\n\n".join(pieces), EXTRACT_MAX_CHARS)
            if len(_tokens(combined)) >= 30:
                return combined

            # OCR fallback for scanned PDFs. It is attempted only when normal extraction is empty.
            try:
                import pytesseract
                from PIL import Image

                ocr_pieces: list[str] = []
                for page_index in range(min(len(pdf), 6)):
                    pix = pdf[page_index].get_pixmap(matrix=fitz.Matrix(1.7, 1.7), alpha=False)
                    image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text = pytesseract.image_to_string(
                        image,
                        lang=os.getenv("TESSERACT_LANG", "eng+ara"),
                        timeout=10,
                    )
                    if text.strip():
                        ocr_pieces.append(f"[Page {page_index + 1} OCR]\n{text.strip()}")
                return _clean("\n\n".join(ocr_pieces), EXTRACT_MAX_CHARS)
            except Exception:
                return combined
    except Exception:
        pass

    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        for page_index, page in enumerate(reader.pages[:max_pages]):
            text = page.extract_text() or ""
            if text.strip():
                pieces.append(f"[Page {page_index + 1}]\n{text.strip()}")
        return _clean("\n\n".join(pieces), EXTRACT_MAX_CHARS)
    except Exception:
        return ""


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document

        doc = Document(str(path))
        pieces: list[str] = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                pieces.append(paragraph.text.strip())
        for table_index, table in enumerate(doc.tables[:40], 1):
            for row in table.rows[:150]:
                values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if values:
                    pieces.append(f"[Table {table_index}] " + " | ".join(values))
        return _clean("\n".join(pieces), EXTRACT_MAX_CHARS)
    except Exception:
        return ""


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        pieces: list[str] = []
        for index, slide in enumerate(prs.slides[:60], 1):
            slide_text: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    slide_text.append(text.strip())
            if slide_text:
                pieces.append(f"[Slide {index}]\n" + "\n".join(slide_text))
        return _clean("\n\n".join(pieces), EXTRACT_MAX_CHARS)
    except Exception:
        return ""


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(path, read_only=True, data_only=True)
        pieces: list[str] = []
        try:
            for sheet in wb.worksheets[:12]:
                pieces.append(f"[Sheet: {sheet.title}]")
                for row_index, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    if row_index > 400:
                        break
                    values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                    if values:
                        pieces.append(" | ".join(values))
                    if sum(len(item) for item in pieces) >= EXTRACT_MAX_CHARS:
                        break
        finally:
            wb.close()
        return _clean("\n".join(pieces), EXTRACT_MAX_CHARS)
    except Exception:
        return ""


def _extract_image(path: Path) -> str:
    try:
        import pytesseract
        from PIL import Image, ImageEnhance, ImageOps

        with Image.open(path) as image:
            image = ImageOps.exif_transpose(image).convert("L")
            image.thumbnail((2600, 2600))
            image = ImageEnhance.Contrast(image).enhance(1.6)
            text = pytesseract.image_to_string(
                image,
                lang=os.getenv("TESSERACT_LANG", "eng+ara"),
                timeout=12,
            )
        return _clean(text, EXTRACT_MAX_CHARS)
    except Exception:
        return ""


def _extract_file(path: Path) -> str:
    suffix = path.suffix.casefold()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix == ".xlsx":
        return _extract_xlsx(path)
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return _extract_image(path)
    if suffix in {".txt", ".md", ".csv"}:
        for encoding in ("utf-8-sig", "utf-8", "cp1256", "cp1252"):
            try:
                return _clean(path.read_text(encoding=encoding, errors="ignore"), EXTRACT_MAX_CHARS)
            except Exception:
                continue
    return ""


def _ground_local(fallback: dict[str, Any], source_context: str, lesson) -> dict[str, Any]:
    result = dict(fallback)
    terms = _distinctive_terms(source_context, 10)
    separator = "، " if lesson.language == "ar" else ", "
    if terms:
        existing = str(result.get("keywords", "")).strip()
        source_terms = separator.join(terms)
        result["keywords"] = f"{existing}{separator if existing else ''}{source_terms}"
    excerpt = next((line.strip() for line in source_context.splitlines() if len(line.strip()) > 35), "")
    if excerpt:
        excerpt = excerpt[:240]
        if lesson.language == "ar":
            result["strategies"] = str(result.get("strategies", "")) + f"\nمرجع الدرس من الملف: {excerpt}"
            result["curriculum"] = str(result.get("curriculum", "")) + "\nتم بناء الخطة من الملف المرفوع بوصفه المرجع الأساسي."
        else:
            result["strategies"] = str(result.get("strategies", "")) + f"\nSource anchor from the uploaded file: {excerpt}"
            result["curriculum"] = str(result.get("curriculum", "")) + "\nThe uploaded file was used as the primary curriculum reference."
    result["_mode"] = "source_grounded_local_fallback"
    return result


def install(core, lesson_engine, lesson_density_patch) -> None:
    if getattr(core, "_lesson_source_grounding_installed", False):
        return

    # Accept the common curriculum file types for lesson preparation as well.
    core.ALLOWED_EXTENSIONS.update({"xlsx", "md", "csv"})

    def source_aware_save_uploaded_file(storage, prefix: str):
        if not storage or not storage.filename:
            return "", ""
        ext = core.file_ext(storage.filename)
        if ext not in core.ALLOWED_EXTENSIONS:
            return storage.filename, f"[Unsupported file type: {ext}]"
        safe_original = secure_filename(storage.filename) or f"source.{ext or 'txt'}"
        saved = Path(core.UPLOAD_DIR) / f"{prefix}_{uuid.uuid4().hex}_{safe_original}"
        storage.save(str(saved))
        extracted = _extract_file(saved)
        if not extracted.strip():
            return storage.filename, (
                "[Unsupported or unreadable file: no readable text was extracted. "
                "For scanned PDF/images, ensure OCR is available or upload a searchable PDF/DOCX.]"
            )
        return storage.filename, extracted

    core.save_uploaded_file = source_aware_save_uploaded_file

    previous_builder = lesson_engine.build_expert_content

    def grounded_builder(lesson, app_module):
        if not str(getattr(lesson, "source_text", "") or "").strip():
            return previous_builder(lesson, app_module)

        context = select_relevant_source(
            lesson.source_text,
            lesson.topic,
            lesson.subject,
            lesson.notes,
            max_chars=max(5000, min(16000, int(os.getenv("LESSON_SOURCE_CONTEXT_CHARS", str(SOURCE_MAX_CHARS))))),
        )
        fallback = lesson_engine.special(lesson, app_module)
        if not context:
            return _ground_local(fallback, lesson.source_text, lesson)

        key = json.dumps(
            [
                SOURCE_CACHE_VERSION,
                lesson.subject,
                lesson.class_name,
                lesson.language,
                lesson.topic,
                lesson.notes,
                hashlib.sha256(context.encode("utf-8", errors="ignore")).hexdigest(),
            ],
            ensure_ascii=False,
        )
        now = time.time()
        cached = lesson_engine.CACHE.get(key)
        if cached and now - cached[0] < lesson_engine.TTL:
            return dict(cached[1])

        if not os.getenv("OPENAI_API_KEY") or lesson_engine.OpenAI is None:
            content = _ground_local(fallback, context, lesson)
            lesson_engine.CACHE[key] = (now, content)
            return dict(content)

        acquired = _AI_GATE.acquire(timeout=max(0.2, float(os.getenv("SOURCE_AI_QUEUE_WAIT", "2.5"))))
        if not acquired:
            content = _ground_local(fallback, context, lesson)
            content["_mode"] = "source_grounded_busy_fallback"
            lesson_engine.CACHE[key] = (now, content)
            return dict(content)

        try:
            target_language = "Arabic" if lesson.language == "ar" else "English"
            strict_system = (
                "You are a senior curriculum specialist producing an official UAE-school lesson plan. "
                "The uploaded source below is the PRIMARY and AUTHORITATIVE reference. Use its actual concepts, vocabulary, examples, text features, and sequence. "
                "Do not replace it with a generic plan and do not import content from another subject. The entered subject is authoritative. "
                "If the subject is English Language, teach the exact reading, writing, grammar, vocabulary, speaking/listening, or literature content present in the source; never reinterpret words such as function, range, table, or figure as mathematics unless the subject itself is Mathematics. "
                "Every major section must be visibly connected to the source. Do not invent quotations, facts, rules, examples, or assessment content not supported by the source. "
                f"Write the plan in {target_language}; retain essential source-language terms and examples where pedagogically necessary. "
                "Return the required structured fields only. Provide exactly six measurable learning outcomes, six student-friendly success criteria, four differentiation levels, and four main activity groups: source-based worked model, guided practice, independent application, and HOTS."
            )
            payload = {
                "subject_entered_by_teacher": lesson.subject,
                "class": lesson.class_name,
                "lesson_title": lesson.topic,
                "periods": lesson.periods,
                "teacher_notes": lesson.notes,
                "source_file": lesson.source_file_name,
                "source_rules": [
                    "Use this source as the primary curriculum reference.",
                    "Use source-specific vocabulary and content in outcomes, activities, assessment, and resources.",
                    "Do not switch subjects because of ambiguous words inside the source.",
                    "When the source does not support a detail, keep the wording cautious rather than inventing it.",
                ],
                "selected_source_context": context,
            }
            client = lesson_engine.OpenAI()
            response = client.responses.parse(
                model=os.getenv("OPENAI_LESSON_MODEL", "gpt-4.1-mini"),
                input=[
                    {"role": "system", "content": strict_system},
                    {"role": "user", "content": json.dumps(payload, ensure_ascii=False)},
                ],
                text_format=lesson_engine.Plan,
                max_output_tokens=max(1800, min(3200, int(os.getenv("OPENAI_LESSON_MAX_TOKENS", "2800")))),
                store=False,
            )
            parsed = response.output_parsed
            data = parsed.model_dump() if parsed else {}
            output: dict[str, Any] = {
                "subject": lesson.subject or fallback.get("subject", ""),
                "class_name": lesson.class_name or fallback.get("class_name", ""),
            }
            for field in lesson_engine.Plan.model_fields:
                output[field] = app_module.clean_text(data.get(field) or fallback.get(field, ""))
            output["learning_outcomes"] = lesson_engine.clean_list(
                output["learning_outcomes"], 6, fallback["learning_outcomes"], lesson.language
            )
            output["success_criteria"] = lesson_engine.clean_list(
                output["success_criteria"], 6, fallback["success_criteria"], lesson.language
            )
            output["differentiation"] = lesson_engine.clean_list(
                output["differentiation"], 4, fallback["differentiation"], lesson.language
            )
            output["main"] = lesson_engine.clean_main(output["main"], fallback["main"], lesson.language)
            output = lesson_density_patch._enrich(lesson, output, fallback)

            # Ensure the resulting plan visibly carries vocabulary from the uploaded source.
            source_terms = _distinctive_terms(context, 10)
            combined = " ".join(str(value) for value in output.values()).casefold()
            overlap = [term for term in source_terms if term.casefold() in combined]
            if source_terms and len(overlap) < min(4, len(source_terms)):
                separator = "، " if lesson.language == "ar" else ", "
                output["keywords"] = separator.join(source_terms) + separator + str(output.get("keywords", ""))
            output["_mode"] = "source_grounded_ai_v5"
            output["_source_file"] = lesson.source_file_name
            lesson_engine.CACHE[key] = (now, output)
            return dict(output)
        except Exception:
            app_module.logger.exception("Source-grounded generation failed; using source-aware local fallback")
            content = _ground_local(fallback, context, lesson)
            content = lesson_density_patch._enrich(lesson, content, fallback)
            content["_mode"] = "source_grounded_error_fallback"
            lesson_engine.CACHE[key] = (now, content)
            return dict(content)
        finally:
            _AI_GATE.release()

    lesson_engine.build_expert_content = grounded_builder

    # Invalidate old Word documents that were generated before uploaded-source grounding.
    try:
        import lesson_concurrency_patch

        lesson_concurrency_patch.CACHE_VERSION = "lesson-docx-source-grounded-v5"
    except Exception:
        pass

    core._lesson_source_grounding_installed = True
